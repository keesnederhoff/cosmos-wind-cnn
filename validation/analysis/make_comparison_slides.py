"""Build a 2-slide deck: (1) SF-Bay wind products compared, (2) the observation
archive collected. Content from the wind-validation synthesis + obs archive.

conda run --live-stream -n dfm_tools python -u make_comparison_slides.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import sys
sys.path.insert(0, str(Path(__file__).resolve().parents[1]))   # project root on path
import config

OUT = config.OUTPUT_ROOT / "SFBay_wind_comparison_slides.pptx"

NAVY  = RGBColor(0x1F, 0x3A, 0x5F)
TEAL  = RGBColor(0x2C, 0x7F, 0x9E)
LIGHT = RGBColor(0xEE, 0xF3, 0xF7)
ALT   = RGBColor(0xFB, 0xFC, 0xFD)
DARK  = RGBColor(0x22, 0x2A, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HILITE = RGBColor(0xB5, 0x44, 0x1A)   # RTMA / "best" accent

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_band_title(slide, title, subtitle):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.15))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
    tf = band.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.45); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = subtitle
    r2.font.size = Pt(13); r2.font.italic = True; r2.font.color.rgb = RGBColor(0xCF, 0xDD, 0xE8)


def style_cell(cell, text, *, bold=False, size=11, color=DARK, fill=WHITE,
               align=PP_ALIGN.LEFT):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color


def add_table(slide, headers, rows, col_w, top, highlight_row0_idx=None):
    nrows, ncols = len(rows) + 1, len(headers)
    total = sum(col_w)
    left = Inches((13.333 - total) / 2.0)
    tbl = slide.shapes.add_table(nrows, ncols, left, Inches(top),
                                 Inches(total), Inches(0.4 * nrows)).table
    tbl.first_row = False; tbl.horz_banding = False
    for j, w in enumerate(col_w):
        tbl.columns[j].width = Inches(w)
    for j, h in enumerate(headers):
        style_cell(tbl.cell(0, j), h, bold=True, size=12, color=WHITE, fill=TEAL)
    for i, row in enumerate(rows, start=1):
        base = LIGHT if i % 2 else ALT
        is_hi = (highlight_row0_idx is not None and (i - 1) == highlight_row0_idx)
        for j, val in enumerate(row):
            style_cell(tbl.cell(i, j), val,
                       bold=is_hi and j == 0, size=11,
                       color=HILITE if is_hi else DARK,
                       fill=RGBColor(0xFC, 0xEF, 0xE6) if is_hi else base)
    return tbl


def sum_to_in(col_w):
    return sum(col_w)


def add_footer(slide, text, top):
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(top), Inches(12.4), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(12); r.font.color.rgb = NAVY


# ---------------- Slide 1: wind products ----------------
s1 = prs.slides.add_slide(BLANK)
add_band_title(s1, "SF Bay Wind Products Compared",
               "10 m wind (+ air temperature) products evaluated for the coupled wave / water-level model")
hdr1 = ["Product", "Type", "Native res.", "Period", "Role / finding"]
rows1 = [
    ["RTMA",       "NOAA 2DVar obs-assimilating analysis", "2.5 km", "2011-present", "Best (obs-constrained); tops wind + thermo"],
    ["HRRR",       "NOAA convection-allowing model",       "3 km",   "2014-present", "Strong 2014+; wins U10"],
    ["WRF_CalNev", "ERA5/WRF downscaling (spd+T only)",    "1.5 km", "1990-2021",    "Best free-running downscaling pre-2011"],
    ["NOW-23",     "ERA5-forced WRF reanalysis",           "2 km",   "2000-2022",    "Near-zero-bias long bridge"],
    ["Sup3rWind",  "ML (GAN / WTK) super-resolution",      "2 km",   "2007-2013",    "Bimodal: good offshore, poor inner-Bay"],
    ["UCLA",       "ERA5-driven WRF",                      "3 km",   "1980-2020",    "Mid-pack; weak U10"],
    ["CONUS404",   "WRF reanalysis",                       "4 km",   "1979-2021",    "Under-spread -> drop for waves"],
    ["ERA5",       "Global reanalysis (baseline)",         "31 km",  "1940-present", "Hard to beat but smooth / under-spread"],
    ["CNN",        "ERA5->CONUS404 ML emulator",           "4 km",   "1940-2027",    "Skill collapses post-2021 extrapolation"],
]
add_table(s1, hdr1, rows1, col_w=[1.5, 3.5, 1.2, 1.6, 4.3], top=1.45, highlight_row0_idx=0)
add_footer(s1,
           "Validated vs the 295-station obs archive (pooled Murphy skill + Taylor; circular RMSE for direction).  "
           "Recommendation: RTMA wherever it exists (2011+); WRF_CalNev best free-running pre-2011; "
           "NOW-23 as the zero-bias bridge. Drop ERA5 / CNN / CONUS404 - under-spread, they suppress the wind "
           "variance waves need.",
           top=6.45)

# ---------------- Slide 2: observations ----------------
s2 = prs.slides.add_slide(BLANK)
add_band_title(s2, "Observation Archive Collected",
               "Free, public SF + North Bay multi-source wind / met / ocean archive (~295 stations, 4 sources)")
hdr2 = ["Source", "Stations", "Period", "Variables"]
rows2 = [
    ["IEM ASOS / AWOS (airports)", "20",        "1914-2026", "air temp, dewpoint, RH, sea-level pressure, wind spd/dir/gust, precip"],
    ["NDBC / NOAA NOS-PORTS",      "33 (24 wind)", "1980-2026", "wind, pressure, air & water temp, waves Hs/Tp/dir (incl. 46026 offshore, 46237 SF Bar)"],
    ["CWOP / APRSWXNET (via MADIS)","251",       "2024-2026", "wind only (citizen sensors; noisiest - QC'd hard)"],
    ["USGS met buoys",             "3",          "campaign",  "nearshore wind + met (Whale's Tail, Emeryville)"],
]
add_table(s2, hdr2, rows2, col_w=[3.0, 1.6, 1.5, 6.0], top=1.45)
add_footer(s2,
           "Hourly CF-NetCDF (station x time), UTC, physical-range gated + per-station acceptance QC (qc_report.csv).  "
           "Deliverables staged under $COSMOS_VALIDATION_DATA_ROOT : per-source NetCDFs, "
           "station-inventory table, bay-wide station "
           "map, 289 per-station timeseries figures.  Built with the resumable, checkpointed pws_scraper toolchain.",
           top=6.35)

prs.save(str(OUT))
print("Wrote", OUT)
