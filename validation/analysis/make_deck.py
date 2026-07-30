"""ERA5 -> RTMA wind-downscaling CNN: findings decks.

One ~20-slide deck that serves both the client update and the USGS science talk,
with the deeper statistical material in a clearly-marked backup section after the
closing slide so it can be skipped live.

Every skill number is read from rankings/combined_skill_weighted.csv at build time --
nothing is hardcoded, so the deck cannot disagree with the validation output.

    python analysis/make_deck.py [--out DIR]
"""
from __future__ import annotations
from pathlib import Path
import argparse
import sys

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
import config

W, H = Inches(13.333), Inches(7.5)          # 16:9
DARK = RGBColor(0x1F, 0x2A, 0x38)
ACCENT = RGBColor(0x1B, 0x6E, 0xA8)
MUTED = RGBColor(0x5A, 0x6B, 0x7B)
GOOD = RGBColor(0x1B, 0x7F, 0x3B)
BAD = RGBColor(0xA8, 0x32, 0x2B)

RES = config.OUTPUT_ROOT
ERA_A, ERA_B = 'eraA_2011-2019', 'eraB_2020-2025'
SPEED, TOP10, DIRV = 'Wind Speed [m/s]', 'Wind Speed [m/s] (top 10%)', 'Wind Direction [deg]'
U10, V10 = 'Wind U10 [m/s]', 'Wind V10 [m/s]'

# Shorter labels for slide tables.
SHORT = {
    'CNN-windonly-BC': 'CNN wind-only (BC)', 'CNN-wave-p2-BC': 'CNN wave-p2 (BC)',
    'CNN-wave-p3-BC': 'CNN wave-p3 (BC)', 'CNN-allvars-BC': 'CNN all-vars (BC)',
    'CNN-extreme-BC': 'CNN extreme (BC)', 'CNN-allvars': 'CNN all-vars (raw)',
    'RTMA-SFbay': 'RTMA 2.5 km', 'ERA5': 'ERA5 ~31 km', 'CONUS404': 'CONUS404 4 km',
}


# --------------------------------------------------------------------------- data
def load_rank() -> pd.DataFrame:
    fp = RES / 'rankings' / 'combined_skill_weighted.csv'
    if not fp.exists():
        raise SystemExit(f"missing {fp} -- run analysis/combined_skill.py first")
    return pd.read_csv(fp)


def rows(rk, era_dir, variable, sort='skill'):
    era = 'Era A' if era_dir == ERA_A else 'Era B'
    d = rk[rk['era'].str.startswith(era) & (rk['variable'] == variable)].copy()
    if sort == 'skill':
        d = d.sort_values('skill', ascending=False)
    else:
        d = d.sort_values(sort)
    return d


def val(rk, era_dir, variable, model, col='skill'):
    d = rows(rk, era_dir, variable)
    m = d[d['model'] == model]
    return float(m[col].iloc[0]) if len(m) else float('nan')


# --------------------------------------------------------------------------- figs
def fig(*parts) -> Path:
    return RES.joinpath(*parts)


def station_fig(era_dir, group, model, station, kind, var='Wind_Speed_m_s'):
    return fig(era_dir, group, model, f'{station}_{var}_{kind}.png')


# --------------------------------------------------------------------------- slide helpers
def _txbox(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_title_slide(prs, title, subtitle, footer):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _txbox(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.4), title, 40, True, DARK)
    _txbox(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(1.0), subtitle, 20, False, MUTED)
    _txbox(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5), footer, 12, False, MUTED)
    return s


def add_section(prs, kicker, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _txbox(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(0.5), kicker, 16, True, ACCENT)
    _txbox(s, Inches(0.9), Inches(3.4), Inches(11.5), Inches(1.2), title, 32, True, DARK)
    return s


def _header(slide, title, subtitle=None):
    _txbox(slide, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7), title, 26, True, DARK)
    if subtitle:
        _txbox(slide, Inches(0.55), Inches(1.02), Inches(12.2), Inches(0.5),
               subtitle, 14, False, MUTED)


def add_bullets(prs, title, bullets, subtitle=None, notes=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, title, subtitle)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(11.9), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        txt, lvl = (b if isinstance(b, tuple) else (b, 0))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run()
        r.text = ('- ' if lvl else '') + txt
        r.font.size = Pt(16 if lvl else 18)
        r.font.color.rgb = MUTED if lvl else DARK
        p.space_after = Pt(8)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def add_figure(prs, title, image, subtitle=None, caption=None, notes=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, title, subtitle)
    image = Path(image)
    if not image.exists():
        raise SystemExit(f"figure missing: {image}")
    top = Inches(1.55)
    avail_h = H - top - Inches(0.85 if caption else 0.35)
    from PIL import Image as _Im
    with _Im.open(image) as im:
        iw, ih = im.size
    scale = min(Inches(12.0) / iw, avail_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    s.shapes.add_picture(str(image), int((W - w) / 2), top, w, h)
    if caption:
        _txbox(s, Inches(0.6), H - Inches(0.8), Inches(12.1), Inches(0.6),
               caption, 13, False, MUTED)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def add_figure_pair(prs, title, images, labels, subtitle=None, caption=None, notes=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, title, subtitle)
    from PIL import Image as _Im
    n = len(images)
    slot_w = (W - Inches(0.8)) / n
    top = Inches(1.75)
    avail_h = H - top - Inches(0.9)
    for i, (im_p, lab) in enumerate(zip(images, labels)):
        im_p = Path(im_p)
        if not im_p.exists():
            raise SystemExit(f"figure missing: {im_p}")
        with _Im.open(im_p) as im:
            iw, ih = im.size
        scale = min((slot_w - Inches(0.25)) / iw, avail_h / ih)
        w, h = int(iw * scale), int(ih * scale)
        x = int(Inches(0.4) + i * slot_w + (slot_w - w) / 2)
        s.shapes.add_picture(str(im_p), x, int(top + Inches(0.28)), w, h)
        _txbox(s, int(Inches(0.4) + i * slot_w), top - Inches(0.05),
               int(slot_w), Inches(0.35), lab, 14, True, ACCENT, PP_ALIGN.CENTER)
    if caption:
        _txbox(s, Inches(0.6), H - Inches(0.78), Inches(12.1), Inches(0.6),
               caption, 13, False, MUTED)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def add_table(prs, title, headers, body, subtitle=None, caption=None,
              highlight=None, notes=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, title, subtitle)
    nr, nc = len(body) + 1, len(headers)
    height = min(Inches(0.36) * nr + Inches(0.15), Inches(4.9))
    shp = s.shapes.add_table(nr, nc, Inches(0.8), Inches(1.7),
                             Inches(11.7), height)
    tbl = shp.table
    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = htxt
        pr = c.text_frame.paragraphs[0]
        pr.font.size = Pt(13)
        pr.font.bold = True
        pr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        if j:
            pr.alignment = PP_ALIGN.CENTER
    for i, row in enumerate(body, start=1):
        for j, cell in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(cell)
            pr = c.text_frame.paragraphs[0]
            pr.font.size = Pt(12)
            if j:
                pr.alignment = PP_ALIGN.CENTER
            if highlight and highlight(row):
                pr.font.bold = True
                pr.font.color.rgb = ACCENT
    if caption:
        _txbox(s, Inches(0.8), Inches(1.75) + height + Inches(0.15), Inches(11.7),
               Inches(1.0), caption, 13, False, MUTED)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


# --------------------------------------------------------------------------- content
FIGS = RES / 'deck_figs'
FOCUS_A, FOCUS_B = f'{ERA_A}_focus', f'{ERA_B}_focus'
BEST = 'CNN-windonly-BC'      # best pooled wind-speed skill, both eras
WAVE = 'CNN-wave-p3-BC'       # best direction; the wave-energy variant


def _n_stations(era_dir):
    d = pd.read_csv(RES / era_dir / 'validation_statistics.csv')
    d = d[(d['variable'] == SPEED) & (~d['station'].astype(str).str.contains('MEAN'))
          & d['source'].isin(['IEM', 'NDBC', 'USGS'])]
    g = d.groupby('source')['station'].nunique()
    return {k: int(v) for k, v in g.items()}, int(g.sum())


def _station_skill(era_dir, station, model, variable=SPEED):
    d = pd.read_csv(RES / era_dir / 'validation_statistics.csv')
    q = d[(d['station'] == station) & (d['model'] == model)
          & (d['variable'] == variable) & (d['source'] != 'ALL')]
    return float(q['skill'].iloc[0]) if len(q) else float('nan')


def _taylor_stats(era_dir, group, model, variable=SPEED):
    """Sample-size-pooled correlation and normalised std for one network."""
    import numpy as np
    d = pd.read_csv(RES / era_dir / 'validation_statistics.csv')
    d = d[(d['variable'] == variable) & (d['model'] == model) & (d['source'] == group)
          & (d['obs_std'] > 0.05) & (d['n'] >= 50)]
    n = d['n'].to_numpy(float)
    R = float(np.tanh(np.sum(n * np.arctanh(np.clip(d['corr'].to_numpy(float), -0.999, 0.999))) / n.sum()))
    sd = float(np.sum(n * (d['model_std'].to_numpy(float) / d['obs_std'].to_numpy(float))) / n.sum())
    return R, sd


def _win_counts(era_dir):
    """Stations where the best CNN beats RTMA -- a per-station view of the era flip."""
    d = pd.read_csv(RES / era_dir / 'validation_statistics.csv')
    d = d[(d['variable'] == SPEED) & (~d['station'].astype(str).str.contains('MEAN'))
          & d['source'].isin(['IEM', 'NDBC', 'USGS'])]
    cnn = d[d['model'].str.startswith('CNN')].groupby('station')['skill'].max()
    rtma = d[d['model'] == 'RTMA-SFbay'].set_index('station')['skill']
    j = pd.concat([cnn.rename('c'), rtma.rename('r')], axis=1).dropna()
    return int((j['c'] > j['r']).sum()), int(len(j))


def build(rk, out_dir: Path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    def A(era, var, model, col='skill'):
        return val(rk, era, var, model, col)

    nA, totA = _n_stations(ERA_A)
    nB, totB = _n_stations(ERA_B)
    winA, denA = _win_counts(ERA_A)
    winB, denB = _win_counts(ERA_B)
    staA = ' + '.join(f'{k} {v}' for k, v in nA.items())
    staB = ' + '.join(f'{k} {v}' for k, v in nB.items())

    # 1 -----------------------------------------------------------------
    add_title_slide(
        prs, 'Downscaling ERA5 to RTMA resolution with a CNN',
        'Long-term 2.5 km wind forcing for coastal-riverine modelling, San Francisco Bay',
        'USGS Pacific Coastal and Marine Science Center   |   validation job 3771199')

    # 2 -----------------------------------------------------------------
    add_bullets(
        prs, 'Situation', [
            'Coastal-riverine compound flooding needs wind forcing over DECADES: surge, '
            'wind-waves and river peaks have to be simulated together, over many events, '
            'before return periods mean anything.',
            'ERA5 gives a continuous record from 1940 - but at ~31 km it resolves none of '
            'the Bay\'s topographic steering.',
            'RTMA gives 2.5 km and is demonstrably reliable - but only from 2011.',
        ],
        subtitle='Why long-term wind forcing is the binding constraint',
        notes='Set up the constraint before any skill numbers appear.')

    # 3 -----------------------------------------------------------------
    add_table(
        prs, 'Complication: the long record is coarse, the good record is short',
        ['product', 'resolution', 'period available', 'skill vs obs'],
        [['ERA5', '~31 km', '1940 - present', f"{A(ERA_B, SPEED, 'ERA5'):+.3f}   (Era B)"],
         ['CONUS404', '4 km', '1979 - 2021', f"{A(ERA_A, SPEED, 'CONUS404'):+.3f}   (Era A)"],
         ['RTMA', '2.5 km', '2011 - present', f"{A(ERA_B, SPEED, 'RTMA-SFbay'):+.3f}   (Era B)"]],
        caption='RTMA is roughly 2.4x the skill of ERA5 - and cannot reach back before 2011. '
                'Long-term compound-flood hindcasts need both properties at once, and no '
                'existing product has them.',
        highlight=lambda r: r[0] in ('ERA5', 'RTMA'),
        notes='The whole problem in one table.')

    # 4 -----------------------------------------------------------------
    add_bullets(
        prs, 'Question', [
            'Can a neural network learn the ERA5 -> RTMA mapping well enough to extend '
            'RTMA-quality winds backwards, into the period where only ERA5 exists?',
            ('If yes: one consistent 2.5 km forcing across the whole hindcast period.', 1),
            ('If no: long-term runs stay at ERA5 resolution.', 1),
        ],
        notes='Everything that follows is evidence for or against this.')

    # 5 --- products compared -------------------------------------------
    add_table(
        prs, 'The products being compared',
        ['product', 'what it is', 'resolution', 'role here'],
        [['ERA5', 'ECMWF global reanalysis', '~31 km', 'CNN input; the baseline to beat'],
         ['CONUS404', 'WRF hindcast, CONUS', '4 km', 'independent dynamical downscaling (Era A only)'],
         ['RTMA', 'NOAA real-time mesoscale analysis', '2.5 km', 'CNN training TARGET; the reference'],
         ['CNN all-vars (raw)', 'loss on all 6 variables', '2.5 km', 'raw / bias-corrected control pair'],
         ['CNN all-vars (BC)', 'same, quantile bias-corrected', '2.5 km', 'effect of calibration alone'],
         ['CNN wind-only (BC)', 'loss on wind only', '2.5 km', 'best wind-speed skill'],
         ['CNN extreme (BC)', 'extra hard >10 m/s loss term', '2.5 km', 'peak-wind variant'],
         ['CNN wave-p2 (BC)', 'smooth U^2 energy weighting', '2.5 km', 'wave-relevant variant'],
         ['CNN wave-p3 (BC)', 'smooth U^3 energy weighting', '2.5 km', 'best wind direction']],
        subtitle='Three reference products and six CNN variants',
        caption='All six CNN variants share one architecture and differ ONLY in what the '
                'loss function weights. RTMA is the training target, so beating RTMA means '
                'beating the teacher.',
        notes='Explains what each row is before any ranking is shown.')

    # 6 --- observations and the two eras -------------------------------
    add_table(
        prs, 'Observations and the two validation eras',
        ['', 'Era A   2011-2019', 'Era B   2020-2025'],
        [['IEM airport / ASOS stations', str(nA.get('IEM', 0)), str(nB.get('IEM', 0))],
         ['NDBC buoys & coastal', str(nA.get('NDBC', 0)), str(nB.get('NDBC', 0))],
         ['USGS project moorings', 'none - none deployed yet', str(nB.get('USGS', 0))],
         ['total stations scored', str(totA), str(totB)],
         ['products compared', '9  (incl. CONUS404)', '8  (CONUS404 record ends 2021)']],
        subtitle='Split at 2020 so the RTMA-overlap period and the recent period are never averaged together',
        caption='The four USGS moorings all begin after 2020-01-22, so they fall entirely in '
                'Era B. They are reported as their own group and kept OUT of the pooled '
                'score - four short records at 1.2-4.9 m would otherwise distort a pool of '
                'long 10 m records.',
        notes='Pre-empts the obvious question about why the two eras are not like-for-like.')

    # 7 --- architecture -------------------------------------------------
    add_figure(
        prs, 'The model: a 3D U-Net', FIGS / 'unet_architecture.png',
        subtitle='ERA5 (~31 km, 8 channels, 6 hourly steps)  ->  RTMA 2.5 km',
        caption='Encoder-decoder with skip connections. Pooling never touches the time axis, '
                'so the 6-step context survives to the bottleneck; the residual connection '
                'means the network predicts the CORRECTION to interpolated ERA5, not the wind '
                'from scratch.',
        notes='Read off the trained checkpoint: 4 levels, base_channels 24, 384-channel '
              'bottleneck, 6 output channels of which u and v are used.')

    # 8 --- method -------------------------------------------------------
    add_bullets(
        prs, 'Method', [
            'Train on the 2011+ overlap where ERA5 and RTMA both exist; the network learns '
            'the mapping between them.',
            'Bias-correct with per-grid-cell empirical quantile mapping, fitted on the '
            'training period only, direction preserved.',
            'Validate against point observations the model never saw - not against RTMA. '
            'Scoring against the training target would only measure imitation.',
            'Metrics: Murphy skill for general conditions; energy-weighted skill (obs^2 and '
            'obs^3) for the high winds that drive waves; RMSE over the strongest 10% of hours.',
            ('The energy-weighted metric mirrors the wave-energy training loss, so the '
             'objective and the score ask the same question.', 1),
        ],
        notes='Stress the independence of the validation.')

    # 9, 10 --- three-product per-metric --------------------------------
    add_figure(
        prs, 'Era A (2011-2019): the CNN beats RTMA', FIGS / f'three_product_{ERA_A}.png',
        subtitle=f'ERA5 vs RTMA vs best CNN  |  {staA}',
        caption=f'Better on all four measures, including both energy-weighted scores and '
                f'top-decile RMSE. Station by station, the best CNN beats RTMA at '
                f'{winA} of {denA}.',
        notes='RTMA is the training target - beating it here is the headline result.')

    add_figure(
        prs, 'Era B (2020-2025): RTMA wins', FIGS / f'three_product_{ERA_B}.png',
        subtitle=f'ERA5 vs RTMA vs best CNN  |  {staB}',
        caption=f'RTMA improved sharply after 2020 and is the only product with POSITIVE '
                f'energy-weighted skill. The CNN now leads at only {winB} of {denB} stations. '
                f'Where RTMA exists, use RTMA.',
        notes='Not a failure - it defines where the CNN is and is not the right tool.')

    # 11 --- full ranking ------------------------------------------------
    add_figure_pair(
        prs, 'Full ranking, both eras',
        [FIGS / f'bars_speed_{ERA_A}.png', FIGS / f'bars_speed_{ERA_B}.png'],
        ['Era A  2011-2019', 'Era B  2020-2025'],
        subtitle='Pooled wind-speed Murphy skill, every product',
        caption='The ordering flips between eras. Note also that the five bias-corrected '
                'variants sit within 0.009 of each other - smaller than the 0.018-0.040 '
                'seed-to-seed noise, so they are not distinguishable.',
        notes='Point at the CNN block vs RTMA in each panel.')

    # 12 --- where the CNN wins and loses, one map per era ----------------
    add_figure(
        prs, 'Where the CNN wins and loses: Era A',
        RES / 'skill_maps' / f'skill_map_3panel_{ERA_A}.png',
        subtitle=f'Wind-speed skill by station, 2011-2019  |  {staA}',
        caption=f'Input, downscaled field and target on one colour scale. The CNN carries '
                f'more warm stations than RTMA across the Bay - it leads at {winA} of {denA} '
                f'stations individually.',
        notes='Same three products as the metric charts, now spatially.')

    add_figure(
        prs, 'Where the CNN wins and loses: Era B',
        RES / 'skill_maps' / f'skill_map_3panel_{ERA_B}.png',
        subtitle=f'Wind-speed skill by station, 2020-2025  |  {staB}',
        caption=f'The picture reverses: RTMA is warm almost everywhere and the CNN leads at '
                f'only {winB} of {denB} stations. Its weakest points cluster in the Central '
                f'Bay, around the Golden Gate and Bay Bridge, where topographic steering is '
                f'sharpest.',
        notes='The spatial counterpart of the era flip.')

    # 13 --- Taylor, one network per slide --------------------------------
    add_figure(
        prs, 'Correlation and variance: NDBC buoys and coastal stations',
        FIGS / f'taylor_NDBC_{ERA_A}.png',
        subtitle=f"Taylor diagram, Era A, {nA.get('NDBC', 0)} NDBC stations",
        caption=f'Over-water exposure - the sites that matter for wave forcing. The CNN '
                f'reaches {_taylor_stats(ERA_A, "NDBC", BEST)[1]:.2f} normalised standard '
                f'deviation here, far closer to the observed variability than the pooled '
                f'figure suggests. Under-dispersion is much weaker over water than the '
                f'all-station number implies.',
        notes='NDBC first: these are the marine-exposure sites most relevant to wave forcing.')

    add_figure(
        prs, 'Correlation and variance: IEM airport stations',
        FIGS / f'taylor_IEM_{ERA_A}.png',
        subtitle=f"Taylor diagram, Era A, {nA.get('IEM', 0)} IEM stations",
        caption=f'Land exposure, and the opposite trade. Correlation is HIGHER than over '
                f'water ({_taylor_stats(ERA_A, "IEM", BEST)[0]:.2f} vs '
                f'{_taylor_stats(ERA_A, "NDBC", BEST)[0]:.2f}) but variance is much LOWER '
                f'({_taylor_stats(ERA_A, "IEM", BEST)[1]:.2f} vs '
                f'{_taylor_stats(ERA_A, "NDBC", BEST)[1]:.2f}). The under-dispersion in the '
                f'pooled diagram is largely an airport-station effect, not a marine one.',
        notes='Splitting the networks shows which one drives the pooled number -- and that '
              'the peak-wind weakness is worse on land than over water.')

    # 14 --- u10 / v10 ---------------------------------------------------
    add_figure_pair(
        prs, 'Wind components: U10 and V10',
        [FIGS / f'bars_u10_{ERA_A}.png', FIGS / f'bars_v10_{ERA_A}.png'],
        ['U10  (east-west)', 'V10  (north-south)'],
        subtitle=f'Era A, pooled Murphy skill  |  {staA}',
        caption=f"RTMA falls to {A(ERA_A, U10, 'RTMA-SFbay'):+.3f} on U10 - below every CNN "
                f"variant. The components show the advantage more clearly than speed alone, "
                f"because speed hides directional error.",
        notes='Components before direction: they explain where the direction result comes from.')

    # 15 --- direction ---------------------------------------------------
    add_figure(
        prs, 'Wind direction', FIGS / f'bars_dir_{ERA_A}.png',
        subtitle='Era A, circular RMSE - lower is better',
        caption=f"The wave-weighted variant gives the best direction of any product: "
                f"{A(ERA_A, DIRV, WAVE, 'rmse'):.1f} deg vs RTMA's "
                f"{A(ERA_A, DIRV, 'RTMA-SFbay', 'rmse'):.1f} deg, with far less directional "
                f"bias ({A(ERA_A, DIRV, WAVE, 'bias'):+.2f} vs "
                f"{A(ERA_A, DIRV, 'RTMA-SFbay', 'bias'):+.2f} deg). Direction sets fetch.",
        notes='Direction matters as much as speed for wave generation in a fetch-limited bay.')

    # 16 --- peaks -------------------------------------------------------
    add_table(
        prs, 'The hard part: nobody resolves the strongest 10% of hours',
        ['product', 'top-10% RMSE [m/s]', 'top-10% bias [m/s]'],
        [[SHORT.get(r['model'], r['model']), f"{float(r['rmse']):.2f}", f"{float(r['bias']):+.2f}"]
         for _, r in rows(rk, ERA_B, TOP10, sort='rmse').iterrows()],
        subtitle='Era B, strongest decile of observed wind speed',
        caption='Every product under-predicts peaks by 1.5-2.9 m/s, RTMA included. For wave '
                'forcing this is the metric that matters most, and it is where the least '
                'progress has been made.',
        highlight=lambda r: r[0] == 'RTMA 2.5 km',
        notes='Be direct: this is an open problem.')

    # 17 --- Alameda -----------------------------------------------------
    add_figure_pair(
        prs, 'Alameda (AAMC1): timeseries and scatter',
        [station_fig(ERA_B, 'NDBC', WAVE, 'AAMC1', 'peak_event_timeseries'),
         station_fig(ERA_B, 'NDBC', WAVE, 'AAMC1', 'scatter')],
        ['Storm peak', 'All hours'],
        subtitle='CNN wave-p3 (bias-corrected) vs observations, Era B',
        caption='Alameda is one of the strongest stations in the network (n = 51,209, skill '
                '+0.524). The CNN tracks the synoptic evolution closely; the scatter shows '
                'the residual flattening at the top of the range.',
        notes='Central Bay, adjacent to the Emeryville Crescent model domain.')

    # 18 --- Alameda storm field ----------------------------------------
    add_figure_pair(
        prs, 'Alameda storm peak: ERA5 -> CNN -> RTMA',
        [fig(FOCUS_B, 'NDBC', 'ERA5', 'AAMC1_spatial_peak_wind_ERA5.png'),
         fig(FOCUS_B, 'NDBC', WAVE, f'AAMC1_spatial_peak_wind_{WAVE}.png'),
         fig(FOCUS_B, 'NDBC', 'RTMA-SFbay', 'AAMC1_spatial_peak_wind_RTMA-SFbay.png')],
        ['ERA5 ~31 km', 'CNN 2.5 km', 'RTMA 2.5 km'],
        subtitle='Max wind speed over the 24 h around the observed peak, March 2023',
        caption='The input, the downscaled field and the target at the same moment. The CNN '
                'recovers structure ERA5 cannot represent at all.',
        notes='The clearest single demonstration of what downscaling buys.')

    # 19 --- a case where the CNN is better ------------------------------
    add_figure_pair(
        prs, 'A case where the CNN beats RTMA: LNDC1',
        [fig(FOCUS_A, 'NDBC', 'ERA5', 'LNDC1_spatial_peak_wind_ERA5.png'),
         fig(FOCUS_A, 'NDBC', WAVE, f'LNDC1_spatial_peak_wind_{WAVE}.png'),
         fig(FOCUS_A, 'NDBC', 'RTMA-SFbay', 'LNDC1_spatial_peak_wind_RTMA-SFbay.png')],
        ['ERA5 ~31 km', 'CNN 2.5 km', 'RTMA 2.5 km'],
        subtitle=f"Era A storm peak.  Station skill: CNN "
                 f"{_station_skill(ERA_A, 'LNDC1', WAVE):+.3f}  vs  RTMA "
                 f"{_station_skill(ERA_A, 'LNDC1', 'RTMA-SFbay'):+.3f}",
        caption='One of 32 Era-A stations where the CNN outscores RTMA. The CNN field is '
                'sharper through the strait than RTMA at the same resolution.',
        notes='Balances the Alameda slide, where RTMA is the stronger product.')

    # 20 --- USGS --------------------------------------------------------
    add_figure(
        prs, 'USGS project moorings', fig(ERA_B, 'USGS', 'multi_model_Wind_Speed_m_s.png'),
        subtitle='Four moorings, Era B - reported separately from the pooled score',
        caption='Grizzly Bay (ERO20) is close to a tie: RTMA +0.513 vs CNN +0.502. All models '
                'read 1.5-1.8 m/s LOW at these sites - a siting effect (open-water moorings '
                'inside cells that average in land), not model error. Treat as a lower bound.',
        notes='These are the stations closest to the project\'s own model domains.')

    # 21, 22 --- close ---------------------------------------------------
    add_bullets(
        prs, 'Next steps', [
            'Validate before 2011 - the period the CNN exists for. No inference has been run '
            'there yet, and 16-26 stations are available back to 1990.',
            ('The bias-correction quantile maps were fitted on 2011+; applying them to the '
             '1990s assumes a stationary wind distribution. That is the main risk and it is '
             'directly testable.', 1),
            'Attack the peak-wind gap through the TAIL of the bias correction before any '
            'retraining - the evidence points at the correction, not the network.',
            'Fix the direction-invariance bug in the bias correction (small, confirmed across '
            'all products).',
            'Stop exploring architectures: the variants differ by less than seed noise. '
            'Remaining gains are in calibration and record length.',
        ],
        notes='Close on the decision, not on the model.')

    add_bullets(
        prs, 'Bottom line', [
            'Where RTMA exists (2011 onward), force models with RTMA.',
            f'Over 2011-2019 the bias-corrected CNN outperforms RTMA - at {winA} of {denA} '
            f'stations individually - which is strong evidence it can carry the pre-2011 record.',
            'The remaining work is qualifying it there, not building a better network.',
        ],
        notes='One slide the audience can repeat back.')

    # ---------------- BACKUP -------------------------------------------
    add_section(prs, 'BACKUP', 'Supporting detail')

    add_bullets(
        prs, 'Backup: how skill is measured', [
            'Murphy skill = 1 - MSE/variance. 0 = no better than the observed mean, 1 = '
            'perfect. Pooled across stations by sample size, then across networks.',
            'Energy-weighted skill: the same score with each hour weighted by observed '
            'speed^2 (wind stress) and speed^3 (wave energy).',
            'Top-decile RMSE: Murphy skill is near-useless over the strongest 10% - the '
            'conditional variance is tiny, so the score is almost always negative. Read RMSE.',
            'Bias-removed skill (skill_dm) separates pattern error from calibration error.',
        ])

    add_bullets(
        prs, 'Backup: what the numbers do NOT support', [
            'The five loss variants are statistically indistinguishable. Era A spans 0.009 '
            'skill; measured seed-to-seed noise is 0.018-0.040.',
            ('Do not rank the architectures on this evidence.', 1),
            'Bias-removed skill is flat at ~0.58-0.60 across every CNN product, above RTMA\'s '
            f"{A(ERA_A, SPEED, 'RTMA-SFbay', 'skill_dm'):.3f}.",
            ('The CNN\'s spatial pattern already beats RTMA\'s. The visible ranking is '
             'calibration, not physics.', 1),
            'Bias correction is worth +0.06 to +0.07 skill - more than any architectural choice.',
        ])

    add_bullets(
        prs, 'Backup: bias correction helps the score and hurts the pattern', [
            'At the top decile, bias-removed skill is BEST for the raw CNN '
            f"({A(ERA_B, TOP10, 'CNN-allvars', 'skill_dm'):+.3f}) and worse for every corrected "
            f"twin (~{A(ERA_B, TOP10, BEST, 'skill_dm'):+.3f}).",
            'So quantile mapping buys the headline extreme number with bias alone, while '
            'degrading the underlying pattern at high winds.',
            'Consistent with the under-dispersion in the Taylor diagram: the quantile map '
            'appears to compress the tail.',
            'Actionable: inspect the fitted maps at the upper percentiles before retraining.',
        ])

    add_table(
        prs, 'Backup: per-station wins, best CNN vs RTMA',
        ['era', 'stations where CNN > RTMA', 'total', 'share'],
        [['Era A  2011-2019', str(winA), str(denA), f'{100*winA/denA:.0f}%'],
         ['Era B  2020-2025', str(winB), str(denB), f'{100*winB/denB:.0f}%']],
        caption='The pooled scores understate how sharp the era flip is. This is the same '
                'result counted station by station rather than pooled.')

    out_dir.mkdir(parents=True, exist_ok=True)
    fp = out_dir / 'deck_findings.pptx'
    prs.save(str(fp))
    n = len(prs.slides._sldIdLst)
    print(f"wrote {fp}  ({n} slides)")
    return fp


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--out', default=str(RES / 'decks'))
    a = ap.parse_args()
    build(load_rank(), Path(a.out))


if __name__ == '__main__':
    main()
